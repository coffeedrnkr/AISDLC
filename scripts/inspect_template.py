
from pptx import Presentation

TEMPLATE_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/Into_CleanDeck_Fixed.pptx"

def inspect_template():
    prs = Presentation(TEMPLATE_FILE)
    
    print(f"--- Layout Inspection for {TEMPLATE_FILE} ---")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\nLayout Index: {i}")
        print(f"  Name: {layout.name}")
        
        # Check placeholders
        for shape in layout.placeholders:
            print(f"    Placeholder idx:{shape.placeholder_format.idx}, type:{shape.placeholder_format.type}, name:'{shape.name}'")
            if hasattr(shape, 'top') and shape.top:
                print(f"      Pos: top={shape.top.inches:.2f}, left={shape.left.inches:.2f}, width={shape.width.inches:.2f}, height={shape.height.inches:.2f}")
            else:
                print("      Pos: (No geometry)")

if __name__ == "__main__":
    inspect_template()
