from pptx import Presentation
import sys

PPTX_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/presentation_deck_generated.pptx"

CONTENT_TO_CHECK = [
    "AI-AUGMENTED SDLC FRAMEWORK",
    "Pillar 1: VS Code as Home Base",
    "Pillar 10: Jira Integration",
    "Bi-directional sync between Markdown artifacts and Jira tickets",
    "Pillar 8: Change Management",
    "Blast Radius",
    "Appendix A: Agent Registry",
    "Appendix B: Prompt Registry"
]

def verify_pptx():
    try:
        prs = Presentation(PPTX_FILE)
    except Exception as e:
        print(f"FAIL: Could not open PPTX file: {e}")
        sys.exit(1)
        
    all_text = []
    
    slide_count = len(prs.slides)
    print(f"Total Slides: {slide_count}")
    
    if slide_count < 10:
        print("FAIL: Slide count is suspiciously low (<10).")
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_text.append(cell.text_frame.text)
        
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_text.append(cell.text_frame.text)
        
        full_slide_text = "\n".join(slide_text)
        all_text.append(full_slide_text)
        # print(f"--- Slide {i+1} ---\n{full_slide_text}\n")


    full_text = "\n".join(all_text)
    
    missing = []
    for item in CONTENT_TO_CHECK:
        if item not in full_text:
            missing.append(item)
            
    if missing:
        print("FAIL: The following content was not found in the PPTX:")
        for m in missing:
            print(f"- {m}")
        sys.exit(1)
    else:
        print("SUCCESS: All checked content was found in the PPTX.")

if __name__ == "__main__":
    verify_pptx()
