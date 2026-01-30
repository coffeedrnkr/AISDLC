
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import re
import os
from pptx.enum.text import PP_ALIGN

INPUT_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/presentation_deck_v2.md"
OUTPUT_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/presentation_deck.pptx"
TEMPLATE_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/Into_CleanDeck_Fixed.pptx"

def create_presentation():
    if os.path.exists(TEMPLATE_FILE):
        print(f"Using template: {TEMPLATE_FILE}")
        prs = Presentation(TEMPLATE_FILE)
    else:
        print("Template not found, using default.")
        prs = Presentation()
    
    # Read Markdown
    with open(INPUT_FILE, 'r') as f:
        lines = [line.rstrip() for line in f.readlines()]

    chunks = [] # List of (type, title, content_list) -> type='title' (H1) or 'slide' (H2)
    agenda_items = [] # Collect H2 titles for Agenda

    current_type = None
    current_title = ""
    current_content = []

    # 1. Parse into chunks based on headers
    # Logic:
    # #  -> H1: New Title Slide Chunk.
    # ## -> H2: New Standard Slide Chunk. Add to Agenda.
    # ### -> H3: Treat as bold text line in current chunk.
    
    for line in lines:
        stripped = line.strip()
        
        if stripped.startswith("# "):
            # Save previous
            if current_type:
                chunks.append((current_type, current_title, current_content))
            
            # Start H1 (Title Slide)
            current_title = stripped[2:].strip()
            current_type = 'title'
            current_content = []
            
        elif stripped.startswith("## "):
            # Save previous
            if current_type:
                chunks.append((current_type, current_title, current_content))
            
            # Start H2 (Standard Slide)
            current_title = stripped[3:].strip()
            current_type = 'slide'
            current_content = []
            agenda_items.append(current_title)
            
        elif stripped.startswith("### "):
            # Treat H3 as Bold Body Text
            # Add a blank line if content exists, then the bold header
            if current_content:
                current_content.append("") 
            current_content.append(f"**{stripped[4:].strip()}**")
            
        else:
            # Regular content, append if we have an active chunk
            if current_type:
                current_content.append(line)
    
    # Append last chunk
    if current_type:
        chunks.append((current_type, current_title, current_content))

    # 2. Process chunks and create slides
    
    # Track if we have added the Agenda yet (after first title slide)
    agenda_added = False
    
    # Helper: Find layout by name substring
    def get_layout(prs, name_keywords):
        # First pass: Look for exact matches or high confidence
        for layout in prs.slide_layouts:
            for kw in name_keywords:
                if kw.lower() in layout.name.lower():
                    # EXCLUDE layouts that sound like Title Slides unless explicitly asking for Title
                    if "title slide" in layout.name.lower() and "title" not in kw.lower():
                        continue
                    return layout
        
        # Fallback: Look for any layout with a Body placeholder (Type 2) and Title (Type 1)
        for layout in prs.slide_layouts:
            has_title = any(ph.placeholder_format.type == 1 for ph in layout.placeholders)
            has_body = any(ph.placeholder_format.type == 2 for ph in layout.placeholders)
            if has_title and has_body:
                return layout
                
        return prs.slide_layouts[0] # Ultimate fallback

    # Helper: Get placeholder by type/idx
    def get_placeholder(slide, idx_list):
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in idx_list:
                return shape
        # Fallback: return any body/object placeholder
        for shape in slide.placeholders:
            if shape.placeholder_format.type in [2, 7]: # BODY=2, OBJECT=7
                return shape
        return None

    for chunk_type, title, content in chunks:
        
        # --- TITLE SLIDE (H1) ---
        if chunk_type == 'title':
            # User request: First slide (Index 0) is Main Title
            layout = prs.slide_layouts[0] 
            slide = prs.slides.add_slide(layout)
            
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # Subtitle?
            if content and len(slide.placeholders) > 1:
                 # Try to find subtitle placeholder (type 4)
                 for ph in slide.placeholders:
                     if ph.placeholder_format.type == 4: # SUBTITLE
                         ph.text = "\n".join([c.strip() for c in content if c.strip()][:3])
                         break
            
            # Force cleanup
            for ph in slide.placeholders:
                if not ph.has_text_frame: continue
                if not ph.text.strip():
                     sp = ph.element
                     parent = sp.getparent()
                     if parent is not None: parent.remove(sp)

            # Generate Agenda immediately after the FIRST Title Slide
            if not agenda_added and agenda_items:
                # For Agenda, use standard Content layout (search for it)
                create_agenda_slide(prs, agenda_items, get_layout(prs, ["Content", "Standard"]))
                agenda_added = True
            continue

        # --- STANDARD SLIDE (H2) ---
        # User request: "3rd is for subsections" -> Index 2 (Section Header)
        # Strategy: Create Divider Slide (Layout 2) FIRST.
        # Then, if there is content, create a Content Slide (Layout 1).
        
        # 1. Create Divider
        divider_layout = prs.slide_layouts[2]
        divider_slide = prs.slides.add_slide(divider_layout)
        if divider_slide.shapes.title:
            divider_slide.shapes.title.text = title
            
        # Parse elements
        elements = parse_content_elements(content)
        if not elements:
            continue # Just a divider, no content
            
        # 2. Prepare for Content
        # Helper to start new CONTENT slide
        def start_content_slide(slide_title, is_table=False):
            if is_table:
                layout = get_layout(prs, ["Table", "Content"]) 
            else:
                # Explicitly look for "Content" or "Standard"
                layout = get_layout(prs, ["Content", "Title and Content", "Standard"])
                
            s = prs.slides.add_slide(layout)
            if s.shapes.title:
                s.shapes.title.text = slide_title
            
            target_ph = get_placeholder(s, [1, 10, 11, 12, 13, 14])
            return s, target_ph

        # Helper to render markdown (bold, code, italic)
        def render_markdown(paragraph, text):
            pattern = r'(\`.*?\`|\*\*.*?\*\*|\*[^*]+?\*)'
            parts = re.split(pattern, text)
            if paragraph.text == text: paragraph.text = ""
            for part in parts:
                if not part: continue
                run = paragraph.add_run()
                if part.startswith("`") and part.endswith("`"):
                    run.text = part[1:-1]; run.font.name = "Courier New"
                elif part.startswith("**") and part.endswith("**"):
                    run.text = part[2:-2]; run.font.bold = True
                elif part.startswith("*") and part.endswith("*") and not part.startswith("**"):
                    run.text = part[1:-1]; run.font.italic = True
                else:
                    run.text = part

        # Start the first content slide immediately
        current_slide, current_ph = start_content_slide(title)
        
        # If we couldn't find a placeholder, fall back to manual textbox (shouldn't happen with this template)
        current_tf = None
        if current_ph and current_ph.has_text_frame:
            current_tf = current_ph.text_frame
            current_tf.clear() # Remove "Click to edit..."
            current_tf.word_wrap = True
        else:
            # Fallback: Create manual textbox so content is NEVER dropped
            print(f"WARNING: No placeholder found on slide '{title}'. Creating manual textbox.")
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5.5)
            textbox = current_slide.shapes.add_textbox(left, top, width, height)
            current_tf = textbox.text_frame
            current_tf.word_wrap = True

        # Determine strict vertical limit (approx 8 items per slide for safety)
        MAX_ITEMS_PER_SLIDE = 8 
        item_count = 0
        
        for elem_type, elem_data in elements:
            
            # --- OVERFLOW CHECK ---
            if item_count >= MAX_ITEMS_PER_SLIDE:
                 current_slide, current_ph = start_content_slide(title + " (Cont.)", is_table=(elem_type=='table'))
                 if current_ph and current_ph.has_text_frame:
                    current_tf = current_ph.text_frame
                    current_tf.clear()
                    current_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                 item_count = 0
            
            # Safety for first slide too
            if current_tf and not current_tf.auto_size:
                 current_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            # --- HANDLE TEXT ---
            if elem_type == 'text':
                print(f"Adding text to slide '{title}': {elem_data[:30]}...")
                
                stripped = elem_data.strip()
                if not stripped: continue
                
                if current_tf:
                    p = current_tf.add_paragraph()
                    p.space_after = Pt(6)
                    
                    is_bullet = stripped.startswith("* ") or stripped.startswith("- ")
                    is_num = re.match(r'^\d+\.', stripped)
                    is_quote = stripped.startswith("> ")
                    is_h3 = stripped.startswith("### ")
                    
                    if is_bullet:
                        p.level = 0
                        processed_text = stripped[2:]
                    elif is_num:
                        processed_text = stripped
                    elif is_quote:
                        p.level = 1 
                        p.font.italic = True
                        processed_text = stripped[2:]
                    elif is_h3:
                        # Treat as sub-header
                        processed_text = stripped[4:]
                        p.font.bold = True
                        p.space_before = Pt(12)
                        p.space_after = Pt(4)
                        p.level = 0 
                    else:
                        processed_text = stripped
                        p.level = 0

                    render_markdown(p, processed_text)
                    item_count += 1
            
            # --- HANDLE TABLE ---
            elif elem_type == 'table':
                print(f"Adding table to slide '{title}'")
                headers = elem_data['headers']
                rows = elem_data['rows']
                
                if current_ph:
                     try:
                        rows_count = len(rows) + 1
                        cols_count = len(headers)
                        
                        try:
                            table_shape = current_ph.insert_table(rows_count, cols_count).table
                        except AttributeError:
                            left = current_ph.left
                            top = current_ph.top
                            width = current_ph.width
                            height = current_ph.height 
                            
                            needed_height = Inches(0.4 * rows_count)
                            if needed_height < height:
                                height = needed_height

                            table_shape = current_slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
                            
                            try:
                                sp = current_ph.element
                                parent = sp.getparent()
                                if parent is not None:
                                    parent.remove(sp)
                            except Exception:
                                pass

                        TABLE_FONT_SIZE = Pt(11)

                        # Header
                        for k, h_text in enumerate(headers):
                            cell = table_shape.cell(0, k)
                            cell.text_frame.word_wrap = True
                            p = cell.text_frame.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER
                            render_markdown(p, h_text.strip())
                            p.font.bold = True
                            p.font.size = TABLE_FONT_SIZE
                        
                        # Rows
                        for r_idx, row_data in enumerate(rows):
                            for c_idx, cell_text in enumerate(row_data):
                                if c_idx < len(headers):
                                    cell = table_shape.cell(r_idx+1, c_idx)
                                    cell.text_frame.word_wrap = True
                                    if cell.text_frame.paragraphs:
                                        p = cell.text_frame.paragraphs[0]
                                        render_markdown(p, cell_text.strip())
                                        p.font.size = TABLE_FONT_SIZE
                        
                        item_count += 4 
                     except Exception as e:
                         pass
        
        # Cleanup ghost text on the slide we just finished
        for ph in current_slide.placeholders:
            if ph.has_text_frame and not ph.text.strip():
                ph.element.getparent().remove(ph.element)
                    
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

def create_agenda_slide(prs, items, layout):
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Agenda"
    
    # Find body placeholder
    tf = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2: # Body
            tf = shape.text_frame
            tf.clear()
            break
            
    if tf:
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.space_after = Pt(10)
    else:
        # Fallback manual
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)
        height = Inches(5.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item

    # Cleanup ghost text on Agenda slide
    for ph in slide.placeholders:
        if ph.has_text_frame and not ph.text.strip():
             sp = ph.element
             parent = sp.getparent()
             if parent is not None: parent.remove(sp)

def parse_content_elements(lines):
    print(f"DEBUG: Parsing {len(lines)} lines of content...")
    elements = []
    table_buffer = []
    in_table = False
    
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            in_table = True
            table_buffer.append(line)
        else:
            if in_table:
                if len(table_buffer) >= 2: 
                    parsed_table = parse_markdown_table(table_buffer)
                    if parsed_table:
                        elements.append(('table', parsed_table))
                    else:
                        for tbl_line in table_buffer:
                             if tbl_line.strip():
                                elements.append(('text', tbl_line.strip()))
                else:
                     for tbl_line in table_buffer:
                         if tbl_line.strip():
                            elements.append(('text', tbl_line.strip()))
                table_buffer = []
                in_table = False
            
            if stripped:
                elements.append(('text', stripped))
                
    if in_table: 
        if len(table_buffer) >= 2:
            parsed_table = parse_markdown_table(table_buffer)
            if parsed_table:
                elements.append(('table', parsed_table))
            else:
                 for tbl_line in table_buffer:
                     if tbl_line.strip():
                        elements.append(('text', tbl_line.strip()))
        else:
             for tbl_line in table_buffer:
                 if tbl_line.strip():
                    elements.append(('text', tbl_line.strip()))
    return elements

def parse_markdown_table(lines):
    headers = [c.strip() for c in lines[0].strip('|').split('|')]
    if not set(lines[1]).issubset(set("|-: ")):
         pass
    rows = []
    for line in lines[2:]:
        row_cols = [c.strip() for c in line.strip('|').split('|')]
        rows.append(row_cols)
    return {'headers': headers, 'rows': rows}

if __name__ == "__main__":
    create_presentation()
