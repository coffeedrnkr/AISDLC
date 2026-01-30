
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
shape = slide.shapes.add_textbox(0, 0, 100, 100)
p = shape.text_frame.add_paragraph()
print(dir(p))
