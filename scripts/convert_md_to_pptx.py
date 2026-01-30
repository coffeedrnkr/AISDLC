import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt

# Constants
MD_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/presentation_deck_v2.md"
# Using the fixed pptx as the template source since potx had issues
TEMPLATE_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/Into_CleanDeck_Fixed.pptx"
OUTPUT_FILE = "/Users/davidarthurs/Projects/AISDLC/00_Introduction/presentation_deck_generated.pptx"

# Layout Indices (based on inspection)
LAYOUT_TITLE_SLIDE = 54  # 2_Title Slide
LAYOUT_CONTENT = 67      # 1_Content_no logo
LAYOUT_TABLE = 64        # 4_Table

MAX_LINES_PER_SLIDE = 12

def parse_markdown(file_path):
    """
    Parses the markdown file into a structured format.
    Returns:
        title: str (Presentation Title)
        slides: list of dicts {'title': str, 'content': list}
            content items can be:
            - {'type': 'text', 'value': str}
            - {'type': 'table', 'header': [], 'rows': []}
    """
    with open(file_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    presentation_title = "Untitled Presentation"
    slides = []
    current_slide = None
    in_table = False
    table_buffer = []

    for line in lines:
        line = line.strip()
        
        # Skip empty lines unless needed, but useful for delimiting
        if not line:
            if in_table:
                # End of table
                if current_slide and table_buffer:
                    current_slide['content'].append(parse_table(table_buffer))
                table_buffer = []
                in_table = False
            continue

        # Presentation Title (H1)
        if line.startswith("# "):
            presentation_title = line[2:].strip()
            continue

        # Slide Title (H2)
        if line.startswith("## "):
            # Close previous table if any
            if in_table and table_buffer:
                if current_slide:
                    current_slide['content'].append(parse_table(table_buffer))
                table_buffer = []
                in_table = False
            
            # Start new slide
            current_slide = {'title': line[3:].strip(), 'content': []}
            slides.append(current_slide)
            continue
        
        # Subheaders (H3, H4) - treat as bold text or sub-sections
        if line.startswith("### ") or line.startswith("#### "):
            if current_slide:
                 current_slide['content'].append({'type': 'text', 'value': f"**{line.lstrip('#').strip()}**"})
            continue

        # Table detection
        if line.startswith("|"):
            in_table = True
            table_buffer.append(line)
            continue
        
        # If we hit non-table line but we were in table
        if in_table:
            if current_slide and table_buffer:
                current_slide['content'].append(parse_table(table_buffer))
            table_buffer = []
            in_table = False

        # Regular text / bullets
        if current_slide:
            # Clean up markdown specific chars if needed, or keep for basic formatting
            current_slide['content'].append({'type': 'text', 'value': line})

    # Flush final table
    if in_table and table_buffer and current_slide:
        current_slide['content'].append(parse_table(table_buffer))

    return presentation_title, slides

def parse_table(table_lines):
    """Parses markdown table lines into structured table data."""
    # simple parser: expects | col | col |
    header = [c.strip() for c in table_lines[0].strip('|').split('|')]
    rows = []
    # Skip separator line (row 1 usually |---|---|)
    start_idx = 1
    if len(table_lines) > 1 and '---' in table_lines[1]:
        start_idx = 2
    
    for line in table_lines[start_idx:]:
        row = [c.strip() for c in line.strip('|').split('|')]
        rows.append(row)
        
    return {'type': 'table', 'header': header, 'rows': rows}

def create_presentation(title, slides):
    prs = Presentation(TEMPLATE_FILE)
    
    # 1. Create Title Slide
    # Note: python-pptx adds new slides at the end. 
    # Since we loaded a template file that might have slides, we might want to clear them or append.
    # The requirement is "convert to pptx", usually implying a fresh deck.
    # Let's remove existing slides to be clean, or just append. 
    # Removing slides in python-pptx is tricky (xml manipulation). 
    # Easier strategy: The template file provided is likely a deck with masters.
    # If it has existing slides, we'll append. Ideally we'd start blank with this template.
    # For now, we append.
    
    slide_layout = prs.slide_layouts[LAYOUT_TITLE_SLIDE]
    slide = prs.slides.add_slide(slide_layout)
    try:
        slide.shapes.title.text = title
    except AttributeError:
        pass # Some layouts might not have a standard title placeholder

    # 2. Create Content Slides
    for slide_data in slides:
        # Check for split need
        # We'll group content into chunks.
        # A table forces a new slide usually to fit well, or checking space.
        # Simple strategy: 
        # - Text chunk -> Slide
        # - Table -> New Slide (Table Layout)
        # - If Text chunk is too long -> Split
        
        # Flatten content to a list of slides-to-be-created
        sub_slides = []
        current_chunk = []
        current_line_count = 0
        
        idx = 0
        while idx < len(slide_data['content']):
            item = slide_data['content'][idx]
            
            if item['type'] == 'table':
                # check for previous text (Intro)
                intro_text = []
                if current_chunk:
                    total_chars = sum(len(x['value']) for x in current_chunk)
                    if total_chars < 1500: 
                        intro_text = current_chunk
                        current_chunk = []
                        current_line_count = 0
                        print(f"DEBUG: Merging Intro Text ({total_chars} chars) with Table")
                    else:
                        sub_slides.append({'type': 'text', 'items': current_chunk})
                        current_chunk = []
                        current_line_count = 0
                        print(f"DEBUG: Intro Text too long ({total_chars} chars), flushing to separate slide")
                
                # Check for next text (Outro)
                outro_text = []
                if idx + 1 < len(slide_data['content']):
                    next_item = slide_data['content'][idx + 1]
                    if next_item['type'] == 'text':
                        if len(next_item['value']) < 1000: 
                             outro_text = [next_item]
                             idx += 1 
                             print(f"DEBUG: Merging Outro Text with Table")

                sub_slides.append({'type': 'table', 'data': item, 'intro_text': intro_text, 'outro_text': outro_text})
            else:
                # Text
                lines = 1 + len(item['value']) // 80 
                if current_line_count + lines > MAX_LINES_PER_SLIDE:
                    sub_slides.append({'type': 'text', 'items': current_chunk})
                    current_chunk = [item]
                    current_line_count = lines
                else:
                    current_chunk.append(item)
                    current_line_count += lines
            
            idx += 1
        
        if current_chunk:
            sub_slides.append({'type': 'text', 'items': current_chunk})

        # Create the actual slides
        base_title = slide_data['title']
        for i, sub in enumerate(sub_slides):
            slide_title = base_title
            if len(sub_slides) > 1:
                slide_title += f" ({i+1}/{len(sub_slides)})"

            if sub['type'] == 'table':
                create_table_slide(prs, slide_title, sub['data'], sub.get('intro_text', []), sub.get('outro_text', []))
            else:
                create_text_slide(prs, slide_title, sub['items'])

    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

def create_text_slide(prs, title, text_items):
    layout = prs.slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)
    
    # Set Title
    if slide.shapes.title:
        slide.shapes.title.text = title
        
    # Find body placeholder
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1: 
            body = shape
            break
    
    if not body:
        for shape in slide.placeholders:
             if shape.has_text_frame: 
                 body = shape
                 break
    
    if body and body.has_text_frame:
        fill_text_frame(body.text_frame, text_items)

def create_table_slide(prs, title, table_data, intro_text=[], outro_text=[]):
    layout = prs.slide_layouts[LAYOUT_TABLE] 
    slide = prs.slides.add_slide(layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Layout Config
    margin_left = Inches(0.5)
    margin_width = Inches(12.0) # Slide width 13.33? Standard wide is 13.33 x 7.5 usually. 
    # Check template size? Assuming standard 16:9 
    
    # Intro
    intro_bottom = Inches(2.0) # Default start of table
    if intro_text:
        text_ph = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 10:
                text_ph = shape
                break
        if text_ph:
            # Resize logic: Give it more height if needed
            # Move Top slightly up if possible? Default is 0.91. 
            # Make height bigger.
            # ERROR FIX: Title ends at ~1.16 (top 0.26 + height 0.90). 
            # 0.91 causes overlap. Moving down to 1.4.
            text_ph.top = Inches(1.4)
            text_ph.height = Inches(1.0) 
            # ERROR FIX: Ensure width/left are set, as they might be defaulting to 0
            text_ph.left = Inches(0.56) # Match template default
            text_ph.width = Inches(12.32)
            
            fill_text_frame(text_ph.text_frame, intro_text)
            intro_bottom = Inches(2.5) # Buffer after text

    # Outro
    outro_top = Inches(6.5) # Default buffer from bottom? Slide height 7.5.
    if outro_text:
        footer_ph = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 12:
                footer_ph = shape
                break
        if footer_ph:
            # Position footer near bottom
            footer_ph.top = Inches(6.0)
            footer_ph.height = Inches(1.0)
            footer_ph.left = Inches(0.56)
            footer_ph.width = Inches(12.32)
            fill_text_frame(footer_ph.text_frame, outro_text)
            outro_top = Inches(5.9) 

    # Table Geometry
    table_top = intro_bottom
    max_table_height = outro_top - table_top - Inches(0.1)
    if max_table_height < Inches(1.0): 
        max_table_height = Inches(1.0) # fallback
        
    rows = len(table_data['rows']) + 1 
    cols = len(table_data['header'])
    
    # Approx height: 0.4 inch per row?
    req_height = Inches(0.4 * rows)
    actual_height = min(req_height, max_table_height)
    
    table = slide.shapes.add_table(rows, cols, margin_left, table_top, margin_width, actual_height).table
    
    # Set Header
    for i, heading in enumerate(table_data['header']):
        cell = table.cell(0, i)
        # Use fill_text_frame to handle bold formatting
        fill_text_frame(cell.text_frame, [{'value': heading}])
        set_cell_style(cell, bold=True, font_size=12)

    # Set Rows
    for r_idx, row_data in enumerate(table_data['rows']):
        for c_idx, cell_data in enumerate(row_data):
            if c_idx < cols: 
                cell = table.cell(r_idx + 1, c_idx)
                fill_text_frame(cell.text_frame, [{'value': cell_data}])
                set_cell_style(cell, bold=False, font_size=10)

def fill_text_frame(text_frame, text_items):
    text_frame.clear() 
    for item in text_items:
        p = text_frame.add_paragraph()
        text = item['value']
        
        parts = re.split(r'(\*\*.*?\*\*)', text)
        for part in parts:
            run = p.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            else:
                run.text = part
        
        if text.startswith('* ') or text.startswith('- '):
            # Clean bullet again
            p.text = "" 
            clean_text = text[2:].strip()
            parts = re.split(r'(\*\*.*?\*\*)', clean_text)
            for part in parts:
                run = p.add_run()
                if part.startswith('**') and part.endswith('**'):
                        run.text = part[2:-2]
                        run.font.bold = True
                else:
                        run.text = part
            p.level = 0
        else:
            p.level = 0

def set_cell_style(cell, bold=False, font_size=10):
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = 'Arial' # or Theme font

def verify_content_integrity(md_path, pptx_path):
    """
    Ensures that every meaningful line of text in the MD exists in the PPTX.
    """
    # 1. Extract significant text from Markdown
    with open(md_path, 'r', encoding='utf-8') as f:
        md_lines = f.readlines()
        
    md_content_items = []
    
    for line in md_lines:
        line = line.strip()
        if not line: continue
        
        # Skip table separators
        if re.match(r'\|?[\s\-:|]+\|?', line) and '---' in line:
            continue
            
        # Clean text
        # Remove Header markers
        clean = line.lstrip("#").strip()
        # Remove table pipes
        clean = clean.replace("|", " ")
        # Remove bullets
        clean = re.sub(r'^[\*\-]\s+', '', clean)
        # Remove bold markers
        clean = clean.replace("**", "")
        # Normalize whitespace
        clean = " ".join(clean.split())
        
        if len(clean) > 3: # Ignore very short artifacts
            md_content_items.append(clean)
            
    # 2. Extract all text from PPTX
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"FAIL: Could not verify PPTX: {e}")
        return

    pptx_text_blob = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptx_text_blob += shape.text + "\n"
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        pptx_text_blob += cell.text_frame.text + "\n"
    
    # Normalize PPTX blob slightly (whitespace)
    pptx_normalized = " ".join(pptx_text_blob.split())
    
    # 3. Verify
    missing_items = []
    for item in md_content_items:
        # Check if this item is in the blob
        # Note: formatting might change spaces, so checking exact substring of "normalized" item in "normalized" blob is safer
        if item not in pptx_normalized:
             # Try partial match? sometimes splitting happens.
             # Strict check first.
             missing_items.append(item)

    if missing_items:
        print("\n" + "="*40)
        print("QUALITY CHECK FAILED: Potential missing text found!")
        print("="*40)
        count = 0
        for m in missing_items:
            # Filter out some noise if needed, but report all for now
            print(f"[MISSING] {m}")
            count += 1
            if count > 20:
                print("... and more ...")
                break
        print("\nPlease inspect manually. Some failures may be due to whitespace/formatting changes.")
    else:
        print("\n" + "="*40)
        print("QUALITY CHECK PASSED: All text appears to be present.")
        print("="*40)

if __name__ == "__main__":
    print("Parsing Markdown...")
    title, slides = parse_markdown(MD_FILE)
    print(f"Found Title: {title}")
    print(f"Found {len(slides)} slides.")
    
    print(f"Presentation saved to {OUTPUT_FILE}")
    print("Generating Presentation...")
    create_presentation(title, slides)
    
    # Run Quality Check
    print("Running Quality Check...")
    verify_content_integrity(MD_FILE, OUTPUT_FILE)
